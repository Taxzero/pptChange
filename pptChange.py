import logging
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')

def map_color(existing_color, color_map):
    try:
        original_hex = str(existing_color).upper()
    except Exception as e:
        logging.error(f"RGBColor 변환 중 오류 발생: {e}")
        return None

    if original_hex in color_map:
        new_hex = color_map[original_hex]
        try:
            r = int(new_hex[0:2], 16)
            g = int(new_hex[2:4], 16)
            b = int(new_hex[4:6], 16)
            return RGBColor(r, g, b)
        except Exception as e:
            logging.error(f"색상 변환 중 오류 발생: {e}")
            return None
    return None

def update_cell_border(cell, color_map):
    try:
        tcPr = cell._tc.get_or_add_tcPr()
    except Exception as e:
        logging.error(f"셀 tc 프로퍼티 접근 실패: {e}")
        return

    border_tags = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    for tag in border_tags:
        for border in tcPr.findall(tag, tcPr.nsmap):
            srgb = border.find("./a:solidFill/a:srgbClr", tcPr.nsmap)
            if srgb is not None:
                current_color = srgb.get("val")
                if current_color and current_color.upper() in color_map:
                    new_color = color_map[current_color.upper()]
                    srgb.set("val", new_color)
                    logging.info(f"표 셀 경계선 색상이 {current_color}에서 {new_color}로 변경되었습니다.")

def update_text_frame(text_frame, font_name, color_map):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            try:
                if run.font.color and run.font.color.rgb is not None:
                    mapped = map_color(run.font.color.rgb, color_map)
                    if mapped is not None:
                        run.font.color.rgb = mapped
            except Exception as e:
                logging.error(f"텍스트 run 색상 처리 중 오류 발생: {e}")

def update_shape(shape, font_name, color_map):
    if hasattr(shape, "shapes") and len(shape.shapes) > 0:
        text_in_group = any(child.has_text_frame for child in shape.shapes)
        for child in shape.shapes:
            update_shape(child, font_name, color_map)
        if text_in_group and hasattr(shape, "line"):
            if not (shape.line.color and hasattr(shape.line.color, "rgb") and shape.line.color.rgb is not None):
                logging.info(f"그룹 도형 '{shape.name}' 내부에 텍스트가 있으나 선 색상이 없어 그룹 윤곽선을 제거합니다.")
                try:
                    shape.line.fill.background()
                except Exception as e:
                    logging.error(f"그룹 윤곽선 제거 중 오류: {e}")
        return

    if shape.has_text_frame:
        update_text_frame(shape.text_frame, font_name, color_map)

    if hasattr(shape, "table"):
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                try:
                    if cell.fill and hasattr(cell.fill, "fore_color") and cell.fill.fore_color.rgb is not None:
                        mapped = map_color(cell.fill.fore_color.rgb, color_map)
                        if mapped is not None:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = mapped
                except Exception as e:
                    logging.error(f"테이블 셀 채우기 색상 업데이트 중 오류: {e}")
                if cell.text_frame:
                    update_text_frame(cell.text_frame, font_name, color_map)
                update_cell_border(cell, color_map)
        return

    if hasattr(shape, "fill"):
        try:
            if shape.fill.fore_color.rgb is not None:
                mapped = map_color(shape.fill.fore_color.rgb, color_map)
                if mapped is not None:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = mapped
        except Exception as e:
            logging.error(f"도형 채우기 색상 업데이트 중 오류: {e}")

    if hasattr(shape, "line"):
        try:
            if "연결선" in shape.name or "화살표" in shape.name:
                logging.info(f"도형 '{shape.name}'은 화살표/연결선으로 판단되어 선 처리를 건너뜁니다.")
            else:
                if not (shape.line.color and hasattr(shape.line.color, "rgb") and shape.line.color.rgb is not None):
                    logging.info(f"도형 '{shape.name}'의 선(윤곽선)에 현재 색상이 없어 윤곽선을 제거합니다.")
                    shape.line.fill.background()
                else:
                    if shape.line.width is not None and shape.line.width > 0:
                        mapped = map_color(shape.line.color.rgb, color_map)
                        if mapped is not None:
                            shape.line.color.rgb = mapped
                            logging.info(f"도형 '{shape.name}'의 선(윤곽선) 색상이 매핑되어 변경되었습니다.")
                        else:
                            logging.info(f"도형 '{shape.name}'의 선(윤곽선) 색상이 매핑되지 않아 기존 색상을 유지합니다.")
                    if shape.line.color and hasattr(shape.line.color, "rgb") and shape.line.color.rgb is not None:
                        if shape.line.color.rgb == RGBColor(255, 255, 255):
                            logging.info(f"도형 '{shape.name}'의 선(윤곽선) 색상이 흰색으로 정의되어 윤곽선을 제거합니다.")
                            shape.line.fill.background()
        except Exception as e:
            logging.error(f"도형 '{shape.name}'의 선 처리 중 오류 발생: {e}")

def apply_global_styles(prs, font_name, color_map):
    for slide in prs.slides:
        for shape in slide.shapes:
            update_shape(shape, font_name, color_map)


try:
    prs = Presentation("1.pptx")
except Exception as e:
    logging.error(f"PPTX 파일 열기 실패: {e}")
    raise

unified_font_name = "맑은 고딕"
color_map = {
    "4F9F9B": "0065B1",
    "C7E3E2": "CBDBEA",
    "46736E": "203864",
    "4D9995": "2F5597",
    "73B9B6": "8FAADC",
    "D1E8E9": "CBDBEA",
    "428683": "1E6FB1",
    "E3F1F0": "DEEBF7",
    "55AAA5": "0070C0",
    "96CEC7": "4F81BD",
    "4B8D89": "1E6FB1",
    "A3D1CF": "9DC3E6",
    "31859C": "4472C4",
    "84C2BF": "9DC3E6"
}

apply_global_styles(prs, unified_font_name, color_map)

try:
    prs.save("2.pptx")
    logging.info("2.pptx 파일이 생성되었습니다.")
except PermissionError as pe:
    logging.error(f"2.pptx 파일 저장 중 오류 발생: {pe}. 파일이 열려있지 않은지 확인하세요.")
except Exception as e:
    logging.error(f"PPTX 파일 저장 중 예기치 않은 오류 발생: {e}")
